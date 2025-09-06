# instagram_reporter.py

import requests
import pandas as pd
import io
import os
import math
from datetime import datetime, timedelta
from typing import Dict, List, Optional
from config import POSTS_PER_SLIDE
import matplotlib.pyplot as plt
import matplotlib.dates as mdates

# PowerPoint Imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_VERTICAL_ANCHOR

#PDF imports
from jinja2 import Environment, FileSystemLoader
from weasyprint import HTML, CSS

class InstagramReporter:
    """
    A class to fetch, analyze, and generate reports for an Instagram Business Account.
    """
    def __init__(self, access_token: str, page_id: str, api_version: str = "v19.0"):
        """
        Initialize the Instagram Reporter.
        """
        self.access_token = access_token
        self.page_id = page_id
        self.base_url = f"https://graph.facebook.com/{api_version}"

    def get_instagram_account_id(self) -> Optional[str]:
        """Get Instagram Business Account ID from the linked Facebook Page ID."""
        url = f"{self.base_url}/{self.page_id}"
        params = {'fields': 'instagram_business_account', 'access_token': self.access_token}
        try:
            response = requests.get(url, params=params)
            response.raise_for_status()
            data = response.json()
            return data.get('instagram_business_account', {}).get('id')
        except requests.RequestException as e:
            print(f"Error getting Instagram account ID: {e}")
            return None

    def get_posts_data(self, start_date: datetime, end_date: datetime) -> List[Dict]:
        """Fetch Instagram posts from the last specified number of days."""
        ig_account_id = self.get_instagram_account_id()
        if not ig_account_id:
            return []
        start_datetime = datetime.combine(start_date, datetime.min.time())
        end_datetime = datetime.combine(end_date, datetime.max.time())
        
        url = f"{self.base_url}/{ig_account_id}/media"
        fields_to_request = (
            'id,caption,media_type,media_url,permalink,timestamp,like_count,comments_count,thumbnail_url,'
            'children{media_url,media_type},'
            'insights.metric(reach,saved,views)' 
        )
        params = {
            'fields': fields_to_request,
            'since': int(start_datetime.timestamp()),
            'until': int(end_datetime.timestamp()),
            'access_token': self.access_token,
            'limit': 100
        }
        
        all_posts = []
        try:
            current_url = url
            while current_url:
                response = requests.get(current_url, params=params if current_url == url else None)
                response.raise_for_status() # This will raise an HTTPError for 4xx/5xx responses
                data = response.json()
                posts = data.get('data', [])
                all_posts.extend(posts)
                current_url = data.get('paging', {}).get('next')
        except requests.exceptions.HTTPError as e:
            # --- NEW, SMARTER ERROR HANDLING ---
            # Inspect the HTTP status code to give a better error message
            status_code = e.response.status_code
            error_message = e.response.json().get('error', {}).get('message', 'No details provided.')
            
            if status_code == 400:
                # A 400 error often means the token is invalid or permissions are wrong
                raise ValueError(
                    f"API Error (Bad Request): The user's access token may be invalid or expired. "
                    f"Please try logging out and in again. (Details: {error_message})"
                )
            elif status_code == 429:
                # API rate limit exceeded
                raise ValueError("The API is currently busy. Please wait a few minutes and try again.")
            else:
                # Generic API error
                raise ValueError(f"An unexpected API error occurred (Status {status_code}). Please try again later.")
        except Exception as e:
            # Catch other errors like network issues
            raise ValueError(f"A network error occurred: {e}")
        

        for post in all_posts:
            if 'insights' in post and 'data' in post['insights']:
                for metric in post['insights']['data']:
                    post[metric['name']] = metric['values'][0]['value']
            post.pop('insights', None)
        return all_posts

    def analyze_posts(self, posts: List[Dict], sort_metric: str) -> Dict:
        """
        Analyze posts performance, segregating by content type (Static vs. Video),
        and generate insights.
        """
        if not posts:
            return {}

        df = pd.DataFrame(posts)
        
        # --- 1. DATA CLEANING & PREPARATION (Same as before) ---
        df['timestamp'] = pd.to_datetime(df['timestamp'])
        df['date'] = df['timestamp'].dt.date
        df['hour'] = df['timestamp'].dt.hour
        df['day_of_week'] = df['timestamp'].dt.day_name()

        for col in ['like_count', 'comments_count', 'reach', 'saved', 'views', 'thumbnail_url']:
            if col not in df.columns: df[col] = 0
        df = df.fillna(0)
        df['total_engagement'] = df['like_count'] + df['comments_count'] + df['saved']
        df['reach'] = df['reach'].replace(0, 1)
        df['engagement_rate_on_reach'] = (df['total_engagement'] / df['reach']).fillna(0) * 100 
        
        # --- 2. SEGREGATE THE DATAFRAME ---
        # Define what we consider 'static' vs 'video' content
        static_types = ['IMAGE', 'CAROUSEL_ALBUM']
        video_types = ['VIDEO'] # Instagram API uses 'VIDEO' for Reels as well
        
        df_static = df[df['media_type'].isin(static_types)]
        df_video = df[df['media_type'].isin(video_types)]
        
        # --- 3. HELPER FUNCTION for repetitive analysis ---
        # This avoids duplicating code and makes our logic cleaner.
        def get_top_bottom_posts(df_subset, metric_to_sort_by):

            if df_subset.empty: return pd.DataFrame(), pd.DataFrame()

            if metric_to_sort_by not in df_subset.columns:
                metric_to_sort_by = 'reach' # Fallback to a safe default

            df_sorted = df_subset.sort_values(metric_to_sort_by, ascending=False)

            num_posts = len(df_sorted)


            #Edge case handling for fewer number of posts to avoid repetition of same post in top and bottom 3    
            if num_posts <= 5: 
                # Ceiling division for the top half to handle odd numbers graciously (e.g., 5 -> 3 top, 2 bottom)
                split_point = (num_posts + 1) // 2
                top_posts = df_sorted.head(split_point)
                bottom_posts = df_sorted.tail(num_posts - split_point)
    
            else: 
                top_posts = df_sorted.head(3)
                bottom_posts = df_sorted.tail(3)

            return top_posts, bottom_posts

        # ANALYZE EACH SEGMENT
        top_static, bottom_static = get_top_bottom_posts(df_static, sort_metric)
        top_video, bottom_video = get_top_bottom_posts(df_video, sort_metric)
        
        columns_to_keep = [
            'id', 'caption', 'media_type', 'permalink', 'timestamp', 'like_count', 
            'comments_count', 'saved', 'views', 'reach', 
            'total_engagement', 'engagement_rate_on_reach', 'media_url', 'thumbnail_url'
        ]
        
        #  COMPILE THE FINAL INSIGHTS DICTIONARY
        
        insights = {
            'all_posts': df[columns_to_keep].to_dict('records'),
            'total_posts': len(df),
            'total_static_posts': len(df_static),
            'total_video_posts': len(df_video),
            
            'avg_engagement_rate': df['engagement_rate_on_reach'].mean(),
            'avg_static_engagement_rate': df_static['engagement_rate_on_reach'].mean() if not df_static.empty else 0,
            'avg_video_engagement_rate': df_video['engagement_rate_on_reach'].mean() if not df_video.empty else 0,

            # Overall stats
            'total_reach': int(df['reach'].sum()),
            'total_views_or_impressions': int(df['views'].sum()),
            'total_likes': int(df['like_count'].sum()),
            'total_comments': int(df['comments_count'].sum()),
            'total_saves': int(df['saved'].sum()),
            
            # Best time can still be calculated on the whole dataset
            'best_posting_hour': df.groupby('hour')['engagement_rate_on_reach'].mean().idxmax(),
            'best_posting_day': df.groupby('day_of_week')['engagement_rate_on_reach'].mean().idxmax(),
            
            # Content type performance is now more explicit
            'content_type_performance': {
                'Static': df_static['engagement_rate_on_reach'].mean() if not df_static.empty else 0,
                'Video': df_video['engagement_rate_on_reach'].mean() if not df_video.empty else 0
            },
            
            # The new, segregated post lists
            'top_3_static': top_static[columns_to_keep].to_dict('records') if not top_static.empty else [],
            'bottom_3_static': bottom_static[columns_to_keep].to_dict('records') if not bottom_static.empty else [],
            'top_3_video': top_video[columns_to_keep].to_dict('records') if not top_video.empty else [],
            'bottom_3_video': bottom_video[columns_to_keep].to_dict('records') if not bottom_video.empty else []
        }
        return insights

    def create_local_csv_report(self, insights: Dict) -> str:
        """
        Creates the CSV report content in-memory and returns it as a string.
        
        """
        # Use an in-memory text buffer instead of a file
        string_buffer = io.StringIO()
        
        # Part 1: Summary Data
        summary_data = [
            ['Metric', 'Value'],
            ['Total Posts', insights.get('total_posts', 0)],
            ['Average Engagement Rate', f"{insights.get('avg_engagement_rate', 0):.2f}%"],
            ['Total Reach', f"{insights.get('total_reach', 0):,}"],
            ['Total Impressions', f"{insights.get('total_impressions', 0):,}"],
            ['Total Likes', f"{insights.get('total_likes', 0):,}"],
            ['Total Comments', f"{insights.get('total_comments', 0):,}"],
            ['Total Saves', f"{insights.get('total_saves', 0):,}"],
            ['Total Video Views', f"{insights.get('total_video_views', 0):,}"],
            ['Best Hour to Post', f"{insights.get('best_posting_hour', 'N/A')}:00"],
            ['Best Day to Post', insights.get('best_posting_day', 'N/A')],
        ]
        # Add content type performance dynamically
        for content_type, performance in insights.get('content_type_performance', {}).items():
            summary_data.append([f'{content_type} Avg Engagement', f"{performance:.2f}%"])
        
        summary_df = pd.DataFrame(summary_data, columns=['Metric', 'Value'])
        
        # --- Part 2: Post Performance Data ---
        post_performance_data = []

        post_categories = {
            'Top Static': insights.get('top_3_static', []),
            'Bottom Static': insights.get('bottom_3_static', []),
            'Top Video': insights.get('top_3_video', []),
            'Bottom Video': insights.get('bottom_3_video', [])
        }
        
        # Process both top and bottom posts
        for category_name, posts in post_categories.items():
            for post in posts:
                post_performance_data.append({
                    'Performance Category': category_name,
                    'Date': pd.to_datetime(post.get('timestamp')).strftime('%Y-%m-%d'),
                    'Media Type': post.get('media_type'),
                    'Reach': post.get('reach', 0),
                    'Views': post.get('views', 0),
                    'Likes': post.get('like_count', 0),
                    'Comments': post.get('comments_count', 0),
                    'Saves': post.get('saved', 0),
                    'Total Engagement': post.get('total_engagement', 0),
                    'Engagement Rate (%)': f"{post.get('engagement_rate_on_reach', 0):.2f}",
                    'Caption': post.get('caption', '')[:200],
                    'Link': post.get('permalink', '')
                })
        
        posts_df = pd.DataFrame(post_performance_data)
        
        # --- Part 3: Write everything to the in-memory buffer ---
        string_buffer.write("INSTAGRAM MONTHLY REPORT\n")
        string_buffer.write(f"Generated on: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n")
        
        string_buffer.write("SUMMARY METRICS\n")
        summary_df.to_csv(string_buffer, index=False)
        
        string_buffer.write("\n\nPOST PERFORMANCE DETAILS\n")
        if not posts_df.empty:
            posts_df.to_csv(string_buffer, index=False)
        else:
            string_buffer.write("No detailed post data to display.\n")
            
        # Return the entire content of the buffer as a single string
        return string_buffer.getvalue()

    def create_powerpoint_report(self, insights: Dict, title_text: str, logo_path: str, sort_metric_display: str) -> io.BytesIO:
        """Creates a PowerPoint presentation."""
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title_text if title_text else "Instagram Performance Report"
        slide.placeholders[1].text = f"Generated on {datetime.now().strftime('%Y-%m-%d')}"
        if logo_path:
            try:
                slide.shapes.add_picture(logo_path, Inches(8.5), Inches(0.5), height=Inches(0.75))
            except FileNotFoundError:
                print(f"⚠️  Logo file not found at '{logo_path}'.")

        summary_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(summary_slide_layout)
        slide.shapes.title.text = "Monthly Performance Summary"
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()

        summary_points = {
        "Total Posts": insights.get('total_posts'),
        f"Static Posts ({insights.get('total_static_posts')}) Avg. Eng Rate": f"{insights.get('avg_static_engagement_rate', 0):.2f}%",
        f"Video Posts ({insights.get('total_video_posts')}) Avg. Eng Rate": f"{insights.get('avg_video_engagement_rate', 0):.2f}%",
        "Total Reach": f"{insights.get('total_reach', 0):,}",
        "Best Day to Post": insights.get('best_posting_day'),
        }
        for key, value in summary_points.items():  # Main summary Slide
            p = tf.add_paragraph(); p.text = f"{key}: {value}"; p.level = 1
        
        self._create_time_series_slide(prs, insights)
        self._create_content_analysis_slide(prs, insights)

        # --- NEW: Four Collage Slides ---
        self._add_collage_slide(prs, insights.get('top_3_static', []), "Top Performing Static Posts", sort_metric_display)
        self._add_collage_slide(prs, insights.get('top_3_video', []), "Top Performing Videos/Reels", sort_metric_display)
        self._add_collage_slide(prs, insights.get('bottom_3_static', []), "Static Posts Needing Improvement", sort_metric_display)
        self._add_collage_slide(prs, insights.get('bottom_3_video', []), "Videos/Reels Needing Improvement", sort_metric_display)

        self._add_annexure_slides(prs, insights) # Adding the annexure

        powerpoint_buffer = io.BytesIO()
        prs.save(powerpoint_buffer)
        powerpoint_buffer.seek(0)

        return powerpoint_buffer

    def _add_collage_slide(self, prs: Presentation, posts: List[Dict], title: str, sort_metric_display: str):
        """Helper function to add a collage slide."""
        if not posts: return
    
        num_posts = len(posts)
        plural_s = "" if num_posts == 1 else "s"
            
        # Replace '3' with the actual number of posts
        dynamic_base = title.replace("3", str(num_posts))
            
        # Construct the final title
        final_title = f"{dynamic_base}{plural_s} (by {sort_metric_display})"

        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = final_title
        positions = [(Inches(0.5), Inches(1.5)), (Inches(3.7), Inches(1.5)), (Inches(6.9), Inches(1.5))]
        for i, post in enumerate(posts):
            if i >= 3: break
            url = post.get('thumbnail_url') if post.get('media_type') == 'VIDEO' else post.get('media_url')
            if not url: continue
            left, top = positions[i]
            try:
                response = requests.get(url, stream=True); response.raise_for_status()
                pic = slide.shapes.add_picture(io.BytesIO(response.content), left, top, width=Inches(2.8))
                
                # Simple Border
                pic.line.color.rgb = RGBColor(220, 220, 220); pic.line.width = Pt(1.5)
                
                # Dynamic textbox positioning
                text_top = pic.top + pic.height + Inches(0.15)
                txBox = slide.shapes.add_textbox(left, text_top, Inches(2.8), Inches(1.5))
                tf = txBox.text_frame; tf.word_wrap = True
                tf.text = (
                    f"Type: {post.get('media_type', 'N/A')}\n"
                    f"Reach: {post.get('reach', 0):,}\n"
                    f"Eng Rate: {post.get('engagement_rate_on_reach', 0):.2f}%"
                )
                for p in tf.paragraphs: p.font.size = Pt(11)
            except Exception as e:
                print(f"❌ Error processing image/text for post {post.get('id')}. Reason: {e}")

    def _add_annexure_slides(self, prs: Presentation, insights: Dict):
        """
        Adds detailed, paginated annexure slides with a clickable link for each post.
        """
        all_posts = insights.get('all_posts', [])
        if not all_posts:
            return

        sorted_posts = sorted(all_posts, key=lambda x: x['timestamp'], reverse=True)
        
        total_pages = math.ceil(len(sorted_posts) / POSTS_PER_SLIDE)

        for page_num, i in enumerate(range(0, len(sorted_posts), POSTS_PER_SLIDE), 1):
            chunk = sorted_posts[i:i + POSTS_PER_SLIDE]
            
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = f"Annexure: All Posts ({page_num}/{total_pages})"

            rows = len(chunk) + 1
            cols = 7
            
            table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(5.5)).table
            
            table.columns[0].width = Inches(1.2)
            table.columns[1].width = Inches(1.5)
            table.columns[2].width = Inches(1.0)
            table.columns[3].width = Inches(1.0)
            table.columns[4].width = Inches(1.0)
            table.columns[5].width = Inches(1.3)
            table.columns[6].width = Inches(1.0)
            
            table.cell(0, 0).text = 'Date'
            table.cell(0, 1).text = 'Type'
            table.cell(0, 2).text = 'Reach'
            table.cell(0, 3).text = 'Views'
            table.cell(0, 4).text = 'Likes'
            table.cell(0, 5).text = 'Eng. Rate'
            table.cell(0, 6).text = 'Link'
            
            for row_idx, post in enumerate(chunk):
               
                row_num = row_idx + 1
                table.cell(row_num, 0).text = pd.to_datetime(post['timestamp']).strftime('%Y-%m-%d')
                table.cell(row_num, 1).text = post['media_type']
                table.cell(row_num, 2).text = f"{post.get('reach', 0):,}"
                table.cell(row_num, 3).text = f"{post.get('views', 0):,}"
                table.cell(row_num, 4).text = f"{post.get('like_count', 0):,}"
                table.cell(row_num, 5).text = f"{post.get('engagement_rate_on_reach', 0):.2f}%"
                
                cell = table.cell(row_num, 6)
                run = cell.text_frame.paragraphs[0].add_run()
                run.text = "View Post"
                run.hyperlink.address = post.get('permalink')
                run.font.color.rgb = RGBColor(12, 95, 204)
                run.font.underline = True

            for row in table.rows:
                for cell in row.cells:
                    # Set vertical alignment for all cells
                    cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                    
                    # Set font properties for each paragraph in the cell
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(14) # Set the font size to 10 points
    
    def create_full_posts_csv(self, insights: Dict) -> str:
        """
        Creates a CSV of ALL posts from the period for a raw data export.
        """
        all_posts = insights.get('all_posts', [])
        if not all_posts:
            return ""

        # Create a DataFrame directly from the list of all posts
        df = pd.DataFrame(all_posts)
        
        # Define and order the columns for the final CSV
        columns_in_order = [
            'timestamp', 'media_type', 'reach', 'views', 'like_count', 
            'comments_count', 'saved', 'total_engagement', 'engagement_rate_on_reach',
            'caption', 'permalink'
        ]
        # Filter the DataFrame to only include these columns in this order
        df = df[columns_in_order]
        
        # Rename columns for better readability in the CSV
        df = df.rename(columns={
            'timestamp': 'Date',
            'media_type': 'Type',
            'reach': 'Reach',
            'views': 'Views',
            'like_count': 'Likes',
            'comments_count': 'Comments',
            'saved': 'Saves',
            'total_engagement': 'Total Engagements',
            'engagement_rate_on_reach': 'Engagement Rate (%)',
            'caption': 'Caption',
            'permalink': 'Link'
        })

        # Format the date column
        df['Date'] = pd.to_datetime(df['Date']).dt.strftime('%Y-%m-%d')
        # Format the Engagement Rate column
        df['Engagement Rate (%)'] = df['Engagement Rate (%)'].apply(lambda x: f"{x:.2f}")

        string_buffer = io.StringIO()
        df.to_csv(string_buffer, index=False)
        return string_buffer.getvalue()
    
    def generate_report(
        self,
        start_date: datetime.date,
        end_date: datetime.date, 
        report_title: str, 
        logo_path: str, 
        sort_metric: str, 
        sort_metric_display: str,
        include_pdf: bool = True,
        include_ppt: bool = True,
        ):
        """
        The main method to generate all reports in-memory.
        Returns:
            A tuple containing the CSV data (as a string) and the PowerPoint data (as a BytesIO object).
        """
        print("📱 Fetching Instagram posts...")
        
        posts = self.get_posts_data(start_date, end_date)
        if not posts:
            # For Streamlit, it's better to raise an error that the app can catch and display.
            raise ValueError("No posts were found for the selected date range. Please try a different range.")

        print(f"✅ Found {len(posts)} posts.")

        print("📊 Analyzing post performance...")

        insights = self.analyze_posts(posts, sort_metric=sort_metric)
        if not insights:
            raise ValueError("Could not generate insights from the fetched data.")

        # --- Generate files in-memory ---
        
        print("\n📝 Creating CSV data in memory...")
        summary_csv_data = self.create_local_csv_report(insights) # Summary Report
        raw_data_csv = self.create_full_posts_csv(insights) # Full report
        pdf_data = self.create_pdf_report(
            insights, 
            report_title, 
            start_date.strftime('%B %d, %Y'), 
            end_date.strftime('%B %d, %Y'),
            sort_metric_display
        ) if include_pdf else None
        
        print("🖼️  Creating PowerPoint presentation in memory...")
        pptx_data = self.create_powerpoint_report(
            insights=insights, 
            title_text=report_title, 
            logo_path=logo_path,
            sort_metric_display=sort_metric_display
        ) if include_ppt else None
        
        print("\n✅ All reports generated successfully in memory.")
        
        # Return the data objects for the Streamlit app to handle
        return summary_csv_data, raw_data_csv, pdf_data, pptx_data



    def _create_time_series_slide(self, prs: Presentation, insights: Dict):
        """Creates a slide with time-series charts: Likes per Day and Reach per Day."""
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Performance Over Time"

        all_posts = insights.get('all_posts', [])
        if not all_posts:
            return

        df = pd.DataFrame(all_posts)
        df['date'] = pd.to_datetime(df['timestamp']).dt.date

        # --- Chart 1: Likes per Day (Left Side) ---
        try:
            likes_per_day = df.groupby('date')['like_count'].sum()
            fig, ax = plt.subplots()
            ax.plot(likes_per_day.index, likes_per_day.values, marker='o', linestyle='-', color='#8884d8')
            ax.set_ylabel('Total Likes')
            ax.set_title('Likes per Day')
            ax.grid(True, linestyle='--', alpha=0.6)
            
            locator = mdates.AutoDateLocator(minticks=4, maxticks=8)
            formatter = mdates.ConciseDateFormatter(locator)
            ax.xaxis.set_major_locator(locator)
            ax.xaxis.set_major_formatter(formatter)
            
            chart_buffer = io.BytesIO()
            plt.savefig(chart_buffer, format='png', bbox_inches='tight')
            plt.close(fig)
            chart_buffer.seek(0)
            slide.shapes.add_picture(chart_buffer, Inches(0.5), Inches(1.5), width=Inches(4.5))
        except Exception as e:
            print(f"⚠️  Could not generate 'Likes per Day' chart. Reason: {e}")

        # --- Chart 2: Reach per Day (Right Side) ---
        try:
            reach_per_day = df.groupby('date')['reach'].sum()
            fig, ax = plt.subplots()
            ax.plot(reach_per_day.index, reach_per_day.values, marker='o', linestyle='-', color='#82ca9d')
            ax.set_ylabel('Total Reach')
            ax.set_title('Reach per Day')
            ax.grid(True, linestyle='--', alpha=0.6)

            locator = mdates.AutoDateLocator(minticks=4, maxticks=8)
            formatter = mdates.ConciseDateFormatter(locator)
            ax.xaxis.set_major_locator(locator)
            ax.xaxis.set_major_formatter(formatter)
            
            chart_buffer = io.BytesIO()
            plt.savefig(chart_buffer, format='png', bbox_inches='tight')
            plt.close(fig)
            chart_buffer.seek(0)
            slide.shapes.add_picture(chart_buffer, Inches(5.5), Inches(1.5), width=Inches(4.5))
        except Exception as e:
            print(f"⚠️  Could not generate 'Reach per Day' chart. Reason: {e}")

    def create_pdf_report(self, insights: Dict, report_title: str, start_date: str, end_date: str, sort_metric_display: str) -> bytes:
        """
        Creates a PDF report in-memory by rendering an HTML template.
        """
        print("--- PDF: Starting PDF generation ---")
        
        # Set up Jinja2 to load the template from the current directory
        env = Environment(loader=FileSystemLoader('.'))
        template = env.get_template("report_template.html")
        
        # Prepare the data context to pass to the template
        # We need to format the dates nicely for display
        context = {
            "report_title": report_title,
            "start_date": start_date,
            "end_date": end_date,
            "insights": insights,
            "sort_metric_display": sort_metric_display
        }
        
        # Render the template with the data
        html_out = template.render(context)
        print("--- PDF: HTML template rendered ---")
        
        # Use WeasyPrint to convert the rendered HTML to a PDF in memory
        pdf_bytes = HTML(string=html_out).write_pdf()
        print("--- PDF: PDF generated successfully in memory ---")
        
        return pdf_bytes
    
    def _create_content_analysis_slide(self, prs: Presentation, insights: Dict):
        """Creates a slide with content analysis charts: Engagement by Type and Format Mix."""
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Content Strategy Analysis"

        # --- Chart 1: Engagement Rate by Content Type (Left Side) ---
        try:
            content_performance = insights.get('content_type_performance', {})
            static_count = insights.get('total_static_posts', 0)
            video_count = insights.get('total_video_posts', 0)
            labels = {'Static': f"Static\n(n={static_count})", 'Video': f"Video\n(n={video_count})"}
            types = list(labels.values())
            rates = [round(r, 2) for r in content_performance.values()]

            fig, ax = plt.subplots()
            bars = ax.bar(types, rates, color=['#8884d8', '#82ca9d'])
            ax.set_ylabel('Average Engagement Rate (%)')
            ax.set_title('Engagement Rate by Content Type')
            ax.set_ylim(0, max(rates) * 1.2 if rates else 1)
            for bar in bars:
                yval = bar.get_height()
                plt.text(bar.get_x() + bar.get_width()/2.0, yval, f'{yval}%', va='bottom', ha='center')
            
            chart_buffer = io.BytesIO()
            plt.savefig(chart_buffer, format='png', bbox_inches='tight')
            plt.close(fig)
            chart_buffer.seek(0)
            slide.shapes.add_picture(chart_buffer, Inches(0.5), Inches(1.5), width=Inches(4.5))
        except Exception as e:
            print(f"⚠️  Could not generate 'Engagement by Type' chart. Reason: {e}")

        # --- Chart 2: Content Format Mix (Right Side) ---
        try:
            all_posts = insights.get('all_posts', [])
            if all_posts:
                df = pd.DataFrame(all_posts)
                format_counts = df['media_type'].value_counts()
                
                fig, ax = plt.subplots()
                ax.pie(format_counts, labels=format_counts.index, autopct='%1.1f%%',
                       startangle=90, colors=['#ffc658', '#00C49F', '#FF8042'])
                ax.axis('equal') # Equal aspect ratio ensures that pie is drawn as a circle.
                ax.set_title('Content Format Mix')
                
                chart_buffer = io.BytesIO()
                plt.savefig(chart_buffer, format='png', bbox_inches='tight')
                plt.close(fig)
                chart_buffer.seek(0)
                slide.shapes.add_picture(chart_buffer, Inches(5.5), Inches(1.5), width=Inches(4.0))
        except Exception as e:
            print(f"⚠️  Could not generate 'Content Format Mix' chart. Reason: {e}")